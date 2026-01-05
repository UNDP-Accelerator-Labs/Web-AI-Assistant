from sys import path as syspath
from os import getcwd, remove, makedirs
from os.path import join, exists, splitext, dirname, basename
from shutil import rmtree
from pandas import to_datetime

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

## LOAD OPERATIONS
syspath.append(join(dirname(__file__), 'operations/'))
from extract import extract as extract_text
syspath.append(join(dirname(__file__), 'vars/'))
from check_MSO_shape import check as check_MSO_shape

def parse_shapes (shapes, download_dir):
	content = []
	## PARSE IMMEDIATE TEXT
	txt_shapes = [
		shp for shp in shapes
		if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
	]
	for shape in txt_shapes:
		if shape.has_text_frame:
			content.append(shape.text)
	## PARSE PLACEHOLDERS
	placeholder_shapes = [
		shp for shp in shapes
		if shp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
	]
	for shape in placeholder_shapes:
		if shape.has_text_frame:
			content.append(shape.text)
	## PARSE AUTOSHAPES
	placeholder_autoshapes = [
		shp for shp in shapes
		if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
	]
	for shape in placeholder_autoshapes:
		if shape.has_text_frame:
			content.append(shape.text)
	## PARESE PICTURES
	picture_shapes = [
		shp for shp in shapes
		if shp.shape_type == MSO_SHAPE_TYPE.PICTURE
	]
	for shape in picture_shapes:
		image = shape.image
		blob = image.blob

		file = join(download_dir, 'image.{}'.format(image.ext))
		with open(file, 'wb') as f:
			f.write(blob)

		content.append('----media/image:\n{}----'.format(extract_text(file)))
		remove(file)

	# ## PARESE MEDIA: TO DO IF NECESSARY
	# media_shapes = [
	# 	shp for shp in shapes
	# 	if shp.shape_type == MSO_SHAPE_TYPE.MEDIA
	# ]
	# print(media_shapes)
	
	## PARSE GROUPS
	group_shapes = [
		shp for shp in shapes
		if shp.shape_type == MSO_SHAPE_TYPE.GROUP
	]
	if len(group_shapes) > 0:
		group_content = []
		for group_shape in group_shapes:
			if len(group_shape.shapes) > 0:
				# print('found group')
				shp_content = parse_shapes(group_shape.shapes, download_dir)
				if len(shp_content) > 0: group_content.append(shp_content)
				# print([check_MSO_shape(shp.shape_type) for shp in group_shape.shapes])
		## UNNEST THE LISTS IN group_content
		group_content = [c for l in group_content for c in l]
		if len(group_content) > 0: 
			for text in group_content:
				if text is not None and text != '': content.append(text)

	return [text for text in content if text is not None and text != '']

def parse (file, datum = {}, complex_content = False):
	pres = Presentation(file)
	
	## HAD TO HARD CODE THESE BECASUE COULD NOT FIND DOC ON HOW TO ITERATE OVER pres.core_properties
	## DOC HERE: https://python-pptx.readthedocs.io/en/latest/dev/analysis/pkg-coreprops.html
	name, ext = splitext(file)
	datum['name'] = basename(file)
	datum['type'] = ext
	datum['title'] = pres.core_properties.title
	datum['subject'] = pres.core_properties.subject
	datum['author'] = pres.core_properties.author
	datum['keywords'] = pres.core_properties.keywords
	try: datum['created'] = to_datetime(pres.core_properties.created).strftime('%m/%d/%Y, %H:%M')
	except: datum['created'] = None
	try: datum['modified'] = to_datetime(pres.core_properties.modified).strftime('%m/%d/%Y, %H:%M')
	except: datum['modified'] = None
	datum['content'] = []

	## ADDITIONAL METADATA
	## FOR THE LIST OF KEYS IN core_properties: https://python-pptx.readthedocs.io/en/latest/api/presentation.html#:~:text=The%20core%20properties%20are%20author,subject%2C%20title%2C%20and%20version.
	datum['additional_metadata'] = {}
	datum['additional_metadata']['comments'] = pres.core_properties.comments
	datum['additional_metadata']['lastModifiedBy'] = pres.core_properties.last_modified_by
	datum['additional_metadata']['revision'] = pres.core_properties.revision
	datum['additional_metadata']['category'] = pres.core_properties.category
	datum['additional_metadata']['content_status'] = pres.core_properties.content_status
	datum['additional_metadata']['identifier'] = pres.core_properties.identifier
	datum['additional_metadata']['language'] = pres.core_properties.language
	datum['additional_metadata']['last_modified_by'] = pres.core_properties.last_modified_by
	datum['additional_metadata']['last_printed'] = pres.core_properties.last_printed
	datum['additional_metadata']['version'] = pres.core_properties.version

	## CREATE A TEMP DIR FOR EXTRACED MEDIA LIKE IMAGES
	download_dir = join(getcwd(), 'ppt_downloads/')
	if not exists(download_dir):
		makedirs(download_dir)

	for slide in pres.slides:
		# print([check_MSO_shape(shp.shape_type) for shp in slide.shapes if shp.shape_type not in [MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.GROUP]])
		datum['content'].append(parse_shapes(slide.shapes, download_dir))
		# print('\n')

	## REMOVE THE TEMP DIR
	rmtree(download_dir)

	datum['content'] = [text for parsed in datum['content'] for text in parsed if text is not None and text != '']
	if complex_content == False: datum['content'] = '\n\n'.join(datum['content'])

	return datum

if __name__ == '__main__':
	file = '/Users/myjyby/Documents/Projects/Teams/scraper/files/Weekly Call February 1st 2022 R&D feedback iteration  .pptx'
	# file = '/Users/myjyby/Documents/Projects/Teams/scraper/files/ACCELERATORLAB_PPT_R4_not_final_but_you_can_use.pptx'
	print(parse(file))